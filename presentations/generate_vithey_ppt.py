"""Generate Vithey App 15-slide pitch deck (.pptx) with ACLEDA branding and app screenshots."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ASSETS_DIR = Path(__file__).resolve().parent / "assets"

# ACLEDA-inspired palette (dark navy + gold from official logo)
DARK_BLUE = RGBColor(0x0B, 0x2D, 0x5B)
LIGHT_BLUE = RGBColor(0x4A, 0x90, 0xD9)
GOLD = RGBColor(0xE8, 0xB9, 0x23)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x21, 0x21, 0x21)
MUTED = RGBColor(0x5C, 0x6B, 0x7A)
FOOTER_TEXT = "Vithey App · ACLEDA Bank AUB App Competition"

MICROSERVICES = [
    "api-gateway", "auth-service", "user-profile-service",
    "content-service", "career-service", "chat-service",
    "finance-service", "file-service", "notification-service", "ai-service",
]

INFRA_COMPONENTS = [
    "Eureka Server — service discovery",
    "Config Server — centralized configuration",
    "PostgreSQL — primary database",
    "Redis — caching layer",
    "RabbitMQ — async messaging",
    "MinIO — object / file storage",
]

FRONTEND_MODULES = [
    "Auth & Onboarding (splash, login, Google, startup)",
    "Media Feed — posters, videos, job cards",
    "Post Detail, Comments, Share, Create Post",
    "Profile — About, Posters, Videos, Jobs, CV",
    "Upload CV & Job Application flow",
    "Finance — verification, invoices, payments",
    "Chat — list, messages, participant profile",
    "AI Chatbot — Vithey assistant with history",
    "Notifications & Settings (11 modules complete)",
]

SLIDES = [
    {
        "type": "title",
        "title": "Vithey App",
        "subtitle": "Team Vithey",
        "tagline": "Organize · Share · Learn · Career",
        "competition": "ACLEDA Bank AUB App Competition 2026",
        "concept": "All-in-one Student Super App for AUB students",
    },
    {
        "type": "team",
        "title": "Team Members",
        "members": [
            ("Moeng Kimheang", "Team Lead / Full-Stack"),
            ("Khorn Molika", "Frontend / UI"),
            ("Heng Liza", "Backend / API"),
            ("Ponleong Bora", "Backend / DevOps"),
            ("Nam Ayheng", "Frontend / Flutter"),
            ("Nao Soksovannarith", "Backend / Database"),
            ("Phon Dyna", "QA / Documentation"),
        ],
    },
    {
        "type": "bullets",
        "title": "Table of Contents",
        "items": [
            "Project Overview & Key Features",
            "System Architecture",
            "Frontend & Backend Tech Stack",
            "Microservices & Infrastructure",
            "Project Plan & DevOps",
            "Security & Deployment",
            "Impact, Demo & Conclusion",
        ],
    },
    {
        "type": "bullets_image",
        "title": "Project Overview",
        "subtitle": "Unified Student Super App",
        "items": [
            "Social Feed — posts, likes, comments, share",
            "Job System — browse jobs, upload CV, apply",
            "Student Finance — verification, fees, invoices",
            "Private Chat — DMs, requests, block/report",
            "AI Assistant — Vithey chatbot with history",
            "Notifications — typed routing to all modules",
            "Settings — account, privacy, security, help",
        ],
        "image": "screen_home.png",
        "image_caption": "Home Feed — social, jobs & media",
    },
    {
        "type": "architecture_image",
        "title": "System Architecture",
        "layers": [
            ("Flutter Client", "iOS & Android mobile app"),
            ("API Gateway", "Single entry · routing · JWT validation"),
            ("Microservices", "10 Spring Boot services + AI service"),
            ("Message Bus", "RabbitMQ events between services"),
            ("Data Layer", "PostgreSQL · Redis · MinIO"),
        ],
        "image": "backend_architecture.png",
    },
    {
        "type": "gallery",
        "title": "App Features (Implemented)",
        "images": [
            ("screen_home.png", "Social Feed"),
            ("screen_finance.png", "Finance"),
            ("screen_chatbot.png", "AI Chatbot"),
            ("screen_profile.png", "Profile & Jobs"),
        ],
    },
    {
        "type": "two_col",
        "title": "Frontend Architecture (Flutter)",
        "left_title": "Tech Stack",
        "left": [
            "GetX — state, routing, DI",
            "Dio + HTTPS — REST API client",
            "flutter_secure_storage — JWT tokens",
            "shared_preferences — theme, language",
            "cached_network_image, intl, file_picker",
            "Modular: 11 feature modules complete",
        ],
        "right_title": "Key Modules",
        "right": FRONTEND_MODULES[:6],
    },
    {
        "type": "microservices",
        "title": "Backend Microservices (Spring Boot 3)",
        "subtitle": "Java 21 · Maven · Clean Architecture · JWT Security",
        "services": MICROSERVICES,
        "infra": INFRA_COMPONENTS[:4],
    },
    {
        "type": "two_col",
        "title": "Infrastructure & DevOps",
        "left_title": "Shared Infrastructure",
        "left": INFRA_COMPONENTS,
        "right_title": "DevOps Pipeline",
        "right": [
            "Git + GitHub — version control",
            "Docker — per-service containerization",
            "docker-compose — local full stack",
            "GitHub Actions — CI per service",
            "Bruno / Postman — API testing",
            "Eureka + Config Server — cloud-ready",
        ],
    },
    {
        "type": "bullets",
        "title": "Project Plan",
        "items": [
            "✅ Phase 1 — Requirement Analysis & prompts",
            "✅ Phase 2 — UI/UX Design (Figma references)",
            "✅ Phase 3 — Backend microservices development",
            "✅ Phase 4 — Flutter frontend (11 modules)",
            "✅ Phase 5 — API integration & mock layer",
            "✅ Phase 6 — Testing (unit + widget tests)",
            "🔲 Phase 7 — Production deployment (GCP/K8s)",
        ],
    },
    {
        "type": "two_col",
        "title": "Security Design",
        "left_title": "Frontend Security",
        "left": [
            "Secure JWT token storage",
            "flutter_secure_storage for credentials",
            "HTTPS-only API communication",
            "Biometric & 2FA placeholders",
            "Privacy toggles — profile, data, tracking",
        ],
        "right_title": "Backend Security",
        "right": [
            "Spring Security + JWT filters",
            "API Gateway token validation",
            "BCrypt password encryption",
            "Role-based access control (RBAC)",
            "Per-service security config",
        ],
    },
    {
        "type": "bullets",
        "title": "Deployment Strategy",
        "items": [
            "Google Cloud Platform — primary target",
            "Kubernetes — container orchestration",
            "OpenShift — enterprise deployment option",
            "Docker images per microservice",
            "Horizontal scaling via Eureka discovery",
            "High availability — stateless services",
        ],
    },
    {
        "type": "bullets",
        "title": "Impact & Benefits for AUB Students",
        "items": [
            "One app replaces fragmented tools",
            "Social networking within campus community",
            "Direct job applications with CV upload",
            "Track tuition & payment invoices in-app",
            "Safe private messaging between students",
            "AI help for finance, jobs & study topics",
            "Built for ACLEDA Bank AUB Competition",
        ],
    },
    {
        "type": "bullets",
        "title": "Conclusion",
        "items": [
            "Full-stack production-ready architecture",
            "10 microservices + Flutter mobile client",
            "All 11 frontend modules implemented",
            "Secure JWT auth end-to-end",
            "Docker + CI/CD deployment ready",
            "Solves real student daily-life problems",
        ],
    },
    {
        "type": "closing",
        "title": "Thank You",
        "subtitle": "Team Vithey",
        "tagline": "Vithey App — Organize · Share · Learn · Career",
        "contact": "Questions welcome",
    },
]


def asset(name: str) -> Path:
    return ASSETS_DIR / name


def set_slide_background(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_header_bar(slide, prs) -> None:
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.1),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = GOLD
    bar.line.fill.background()


def add_footer(slide, prs, page_num: int) -> None:
    footer_h = Inches(0.32)
    footer = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0), prs.slide_height - footer_h, prs.slide_width, footer_h,
    )
    footer.fill.solid()
    footer.fill.fore_color.rgb = DARK_BLUE
    footer.line.fill.background()

    left = slide.shapes.add_textbox(Inches(0.35), prs.slide_height - Inches(0.29), Inches(9), Inches(0.22))
    p = left.text_frame.paragraphs[0]
    p.text = FOOTER_TEXT
    p.font.size = Pt(9)
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"

    right = slide.shapes.add_textbox(prs.slide_width - Inches(0.9), prs.slide_height - Inches(0.29), Inches(0.6), Inches(0.22))
    p = right.text_frame.paragraphs[0]
    p.text = str(page_num)
    p.font.size = Pt(9)
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.RIGHT


def add_slide_image(slide, filename: str, left, top, width, height=None) -> None:
    path = asset(filename)
    if not path.exists():
        return
    if height:
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))
    else:
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width))


def add_title_text(slide, text: str, top: float, left: float = 0.55, width: float = 12, size: int = 30) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.75))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = "Calibri"


def add_subtitle_text(slide, text: str, top: float, left: float = 0.55, color=LIGHT_BLUE) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(12), Inches(0.45))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(16)
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.font.bold = True


def add_bullet_list(slide, items: list[str], left: float, top: float, width: float, font_size: int = 17) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK_TEXT
        p.font.name = "Calibri"
        p.space_after = Pt(5)
        p.bullet = True


def add_two_columns(slide, left_title, left_items, right_title, right_items) -> None:
    def col(x, title, items):
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.45), Inches(5.9), Inches(5.0),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA)
        card.line.color.rgb = LIGHT_BLUE

        tb = slide.shapes.add_textbox(Inches(x + 0.2), Inches(1.58), Inches(5.5), Inches(0.35))
        p = tb.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE

        bb = slide.shapes.add_textbox(Inches(x + 0.3), Inches(2.0), Inches(5.4), Inches(4.2))
        tf = bb.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item
            p.font.size = Pt(14)
            p.font.color.rgb = DARK_TEXT
            p.space_after = Pt(4)
            p.bullet = True

    col(0.45, left_title, left_items)
    col(6.75, right_title, right_items)


def build_title_slide(slide, prs, data, page_num):
    set_slide_background(slide, WHITE)
    add_header_bar(slide, prs)

    # Left: ACLEDA Bank logo
    add_slide_image(slide, "acleda_bank_logo.png", left=0.4, top=0.35, width=3.2)

    # Right side hero panel
    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(4.0), Inches(0.1), Inches(9.33), Inches(4.5),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = DARK_BLUE
    panel.line.fill.background()

    # Vithey app logo inside panel
    add_slide_image(slide, "vithey_logo.png", left=8.8, top=0.55, width=1.4)

    # Title text in panel
    tb = slide.shapes.add_textbox(Inches(4.4), Inches(1.2), Inches(8.5), Inches(1.0))
    p = tb.text_frame.paragraphs[0]
    p.text = data["title"]
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = WHITE

    sb = slide.shapes.add_textbox(Inches(4.4), Inches(2.2), Inches(8.5), Inches(0.5))
    p = sb.text_frame.paragraphs[0]
    p.text = data["subtitle"]
    p.font.size = Pt(24)
    p.font.color.rgb = GOLD

    tg = slide.shapes.add_textbox(Inches(4.4), Inches(2.85), Inches(8.5), Inches(0.45))
    p = tg.text_frame.paragraphs[0]
    p.text = data["tagline"]
    p.font.size = Pt(20)
    p.font.color.rgb = LIGHT_BLUE
    p.font.bold = True

    # Competition line below panel
    comp = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(4.75), Inches(12.5), Inches(0.9),
    )
    comp.fill.solid()
    comp.fill.fore_color.rgb = RGBColor(0xFF, 0xF8, 0xE8)
    comp.line.color.rgb = GOLD

    cb = slide.shapes.add_textbox(Inches(0.6), Inches(4.88), Inches(12.2), Inches(0.7))
    tf = cb.text_frame
    tf.paragraphs[0].text = data.get("competition", "")
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = DARK_BLUE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = data.get("concept", "")
    p2.font.size = Pt(14)
    p2.font.color.rgb = MUTED
    p2.alignment = PP_ALIGN.CENTER

    # Feature pills
    features = ["Social", "Jobs & CV", "Finance", "Chat", "AI Assistant"]
    x = 0.55
    for feat in features:
        pill = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(5.85), Inches(2.35), Inches(0.45),
        )
        pill.fill.solid()
        pill.fill.fore_color.rgb = LIGHT_BLUE
        pill.line.fill.background()
        pb = slide.shapes.add_textbox(Inches(x), Inches(5.92), Inches(2.35), Inches(0.35))
        p = pb.text_frame.paragraphs[0]
        p.text = feat
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        x += 2.5

    add_footer(slide, prs, page_num)


def build_team_slide(slide, prs, data, page_num):
    set_slide_background(slide, WHITE)
    add_header_bar(slide, prs)
    add_title_text(slide, data["title"], top=0.3)

    # Team table header
    headers = ["Name", "Role"]
    y = 1.2
    for i, h in enumerate(headers):
        hb = slide.shapes.add_textbox(Inches(0.6 + i * 6.2), Inches(y), Inches(5.8), Inches(0.35))
        p = hb.text_frame.paragraphs[0]
        p.text = h
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        bg = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55 + i * 6.2), Inches(y), Inches(5.9), Inches(0.38),
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = DARK_BLUE
        bg.line.fill.background()
        slide.shapes._spTree.remove(hb._element)
        slide.shapes._spTree.insert(2 + i, hb._element)

    for row_i, (name, role) in enumerate(data["members"]):
        ry = 1.65 + row_i * 0.52
        bg_color = RGBColor(0xF5, 0xF7, 0xFA) if row_i % 2 == 0 else WHITE
        row_bg = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(ry), Inches(12.2), Inches(0.48),
        )
        row_bg.fill.solid()
        row_bg.fill.fore_color.rgb = bg_color
        row_bg.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)

        nb = slide.shapes.add_textbox(Inches(0.75), Inches(ry + 0.08), Inches(5.5), Inches(0.35))
        p = nb.text_frame.paragraphs[0]
        p.text = name
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE

        rb = slide.shapes.add_textbox(Inches(6.95), Inches(ry + 0.08), Inches(5.5), Inches(0.35))
        p = rb.text_frame.paragraphs[0]
        p.text = role
        p.font.size = Pt(14)
        p.font.color.rgb = MUTED

    # Team size badge
    badge = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, Inches(11.0), Inches(0.25), Inches(1.8), Inches(0.75),
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = GOLD
    badge.line.fill.background()
    bb = slide.shapes.add_textbox(Inches(11.0), Inches(0.42), Inches(1.8), Inches(0.4))
    p = bb.text_frame.paragraphs[0]
    p.text = "7 Members"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

    add_footer(slide, prs, page_num)


def build_bullets_image_slide(slide, prs, data, page_num):
    set_slide_background(slide, WHITE)
    add_header_bar(slide, prs)
    add_title_text(slide, data["title"], top=0.3, width=6)
    if data.get("subtitle"):
        add_subtitle_text(slide, data["subtitle"], top=0.85)
    add_bullet_list(slide, data["items"], left=0.55, top=1.35, width=6.2, font_size=15)

    # Phone mockup frame
    frame = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(7.1), Inches(0.95), Inches(3.2), Inches(5.6),
    )
    frame.fill.solid()
    frame.fill.fore_color.rgb = DARK_BLUE
    frame.line.color.rgb = GOLD

    add_slide_image(slide, data["image"], left=7.25, top=1.1, width=2.9, height=5.3)

    if data.get("image_caption"):
        cap = slide.shapes.add_textbox(Inches(7.1), Inches(6.6), Inches(3.2), Inches(0.3))
        p = cap.text_frame.paragraphs[0]
        p.text = data["image_caption"]
        p.font.size = Pt(11)
        p.font.color.rgb = MUTED
        p.alignment = PP_ALIGN.CENTER

    add_footer(slide, prs, page_num)


def build_architecture_image_slide(slide, prs, data, page_num):
    set_slide_background(slide, WHITE)
    add_header_bar(slide, prs)
    add_title_text(slide, data["title"], top=0.3)

    # Left: layer stack
    y = 1.0
    for i, (layer, desc) in enumerate(data["layers"]):
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(y), Inches(5.8), Inches(0.62),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xE8, 0xEE, 0xF7) if i % 2 == 0 else RGBColor(0xFF, 0xF8, 0xE8)
        shape.line.color.rgb = DARK_BLUE
        tf = shape.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = f"{layer}  →  {desc}"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE
        y += 0.72

    # Right: architecture diagram
    add_slide_image(slide, data["image"], left=6.6, top=0.85, width=6.2, height=5.5)
    add_footer(slide, prs, page_num)


def build_gallery_slide(slide, prs, data, page_num):
    set_slide_background(slide, WHITE)
    add_header_bar(slide, prs)
    add_title_text(slide, data["title"], top=0.3)

    images = data["images"]
    positions = [(0.55, 1.1), (3.45, 1.1), (6.35, 1.1), (9.25, 1.1)]
    for (img, label), (x, y) in zip(images, positions):
        frame = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.7), Inches(4.9),
        )
        frame.fill.solid()
        frame.fill.fore_color.rgb = DARK_BLUE
        frame.line.color.rgb = GOLD
        add_slide_image(slide, img, left=x + 0.1, top=y + 0.12, width=2.5, height=4.35)
        lb = slide.shapes.add_textbox(Inches(x), Inches(y + 5.05), Inches(2.7), Inches(0.3))
        p = lb.text_frame.paragraphs[0]
        p.text = label
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE
        p.alignment = PP_ALIGN.CENTER

    add_footer(slide, prs, page_num)


def build_microservices_slide(slide, prs, data, page_num):
    set_slide_background(slide, WHITE)
    add_header_bar(slide, prs)
    add_title_text(slide, data["title"], top=0.3)
    if data.get("subtitle"):
        add_subtitle_text(slide, data["subtitle"], top=0.82)

    services = data["services"]
    cols = 2
    per_col = (len(services) + 1) // 2
    for i, svc in enumerate(services):
        col = i // per_col
        row = i % per_col
        x = 0.55 + col * 6.3
        y = 1.35 + row * 0.55
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(5.9), Inches(0.46),
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xE8, 0xF4, 0xFD)
        box.line.color.rgb = LIGHT_BLUE
        tb = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.08), Inches(5.6), Inches(0.32))
        p = tb.text_frame.paragraphs[0]
        p.text = f"●  {svc}"
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_BLUE

    # Infra strip at bottom
    strip = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(5.5), Inches(12.2), Inches(1.2),
    )
    strip.fill.solid()
    strip.fill.fore_color.rgb = RGBColor(0xFF, 0xF8, 0xE8)
    strip.line.color.rgb = GOLD
    tb = slide.shapes.add_textbox(Inches(0.75), Inches(5.62), Inches(11.8), Inches(1.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Shared: " + "  ·  ".join(data.get("infra", []))
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_TEXT

    add_footer(slide, prs, page_num)


def build_closing_slide(slide, prs, data, page_num):
    set_slide_background(slide, DARK_BLUE)
    add_header_bar(slide, prs)

    add_slide_image(slide, "acleda_bank_logo.png", left=0.5, top=0.4, width=2.8)
    add_slide_image(slide, "vithey_logo.png", left=10.5, top=0.5, width=1.5)

    tb = slide.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(11.3), Inches(1.0))
    p = tb.text_frame.paragraphs[0]
    p.text = data["title"]
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    sb = slide.shapes.add_textbox(Inches(1.0), Inches(3.35), Inches(11.3), Inches(0.55))
    p = sb.text_frame.paragraphs[0]
    p.text = data["subtitle"]
    p.font.size = Pt(28)
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER

    tg = slide.shapes.add_textbox(Inches(1.0), Inches(4.05), Inches(11.3), Inches(0.45))
    p = tg.text_frame.paragraphs[0]
    p.text = data["tagline"]
    p.font.size = Pt(18)
    p.font.color.rgb = LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER

    add_footer(slide, prs, page_num)


def build_content_slide(slide, prs, data, page_num):
    set_slide_background(slide, WHITE)
    add_header_bar(slide, prs)
    add_title_text(slide, data["title"], top=0.3)

    slide_type = data["type"]
    if slide_type == "bullets":
        if data.get("subtitle"):
            add_subtitle_text(slide, data["subtitle"], top=0.82)
            add_bullet_list(slide, data["items"], left=0.55, top=1.3, width=12)
        else:
            add_bullet_list(slide, data["items"], left=0.55, top=1.0, width=12)
    elif slide_type == "two_col":
        add_two_columns(slide, data["left_title"], data["left"], data["right_title"], data["right"])

    add_footer(slide, prs, page_num)


def generate(output_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    builders = {
        "title": build_title_slide,
        "team": build_team_slide,
        "bullets_image": build_bullets_image_slide,
        "architecture_image": build_architecture_image_slide,
        "gallery": build_gallery_slide,
        "microservices": build_microservices_slide,
        "closing": build_closing_slide,
    }

    for i, data in enumerate(SLIDES, start=1):
        slide = prs.slides.add_slide(blank)
        slide_type = data["type"]
        builder = builders.get(slide_type, build_content_slide)
        builder(slide, prs, data, i)

    prs.save(str(output_path))
    print(f"Created: {output_path}")
    print(f"Slides: {len(SLIDES)}")
    print(f"Assets: {ASSETS_DIR}")


if __name__ == "__main__":
    out = Path(__file__).resolve().parent / "Vithey_App_Presentation.pptx"
    generate(out)
