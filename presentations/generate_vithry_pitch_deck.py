"""
Vithry App — Professional 15-slide pitch deck generator.

Installation:
    pip install python-pptx pillow

Usage:
    python generate_vithry_pitch_deck.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Brand palette
# ---------------------------------------------------------------------------
DARK_BLUE = RGBColor(0x0B, 0x2D, 0x5B)
LIGHT_BLUE = RGBColor(0x4A, 0x90, 0xD9)
GOLD = RGBColor(0xE8, 0xB9, 0x23)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x21, 0x21, 0x21)
MUTED = RGBColor(0x6B, 0x7B, 0x8C)
LIGHT_BG = RGBColor(0xF4, 0xF7, 0xFB)
CARD_BG = RGBColor(0xE8, 0xF1, 0xFB)
PLACEHOLDER_BG = RGBColor(0xE0, 0xE4, 0xEA)

FOOTER_LEFT = "Vithry App"
OUTPUT = Path(__file__).resolve().parent / "Vithry_App_Pitch_Deck.pptx"

TEAM = [
    "Moeng Kimheang",
    "Khorn Molika",
    "Heng Liza",
    "Ponleong Bora",
    "Nam Ayheng",
    "Nao Soksovannarith",
    "Phon Dyna",
]

TOC_ITEMS = [
    "Project Overview",
    "System Architecture",
    "Tech Stack",
    "Infrastructure",
    "Project Plan",
    "DevOps Strategy",
    "Backend Design",
    "Frontend Design",
    "Security Design",
    "Deployment Strategy",
    "Conclusion",
]


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def new_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_top_accent(slide, prs) -> None:
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.09),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = GOLD
    bar.line.fill.background()


def add_footer(slide, prs, page: int) -> None:
    h = Inches(0.34)
    y = prs.slide_height - h
    strip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), y, prs.slide_width, h)
    strip.fill.solid()
    strip.fill.fore_color.rgb = DARK_BLUE
    strip.line.fill.background()

    left = slide.shapes.add_textbox(Inches(0.4), y + Inches(0.06), Inches(4), Inches(0.22))
    p = left.text_frame.paragraphs[0]
    p.text = FOOTER_LEFT
    p.font.name = "Calibri"
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE

    right = slide.shapes.add_textbox(prs.slide_width - Inches(0.9), y + Inches(0.06), Inches(0.55), Inches(0.22))
    p = right.text_frame.paragraphs[0]
    p.text = str(page)
    p.font.name = "Calibri"
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.RIGHT


def add_title(slide, text: str, top: float = 0.22, size: int = 32, color=DARK_BLUE, center=False) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(top), Inches(12.2), Inches(0.7))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.name = "Calibri"
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = color
    if center:
        p.alignment = PP_ALIGN.CENTER


def add_subtitle(slide, text: str, top: float, center=False, color=LIGHT_BLUE) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(top), Inches(12.2), Inches(0.45))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.name = "Calibri"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = color
    if center:
        p.alignment = PP_ALIGN.CENTER


def add_bullets(slide, items: list[str], left: float, top: float, width: float, size: int = 17) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.bullet = True
        p.font.name = "Calibri"
        p.font.size = Pt(size)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(6)


def add_image_placeholder(slide, left: float, top: float, width: float, height: float, label: str = "Image Placeholder") -> None:
    rect = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height),
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = PLACEHOLDER_BG
    rect.line.color.rgb = LIGHT_BLUE
    rect.line.width = Pt(1.5)

    icon = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, Inches(left + width / 2 - 0.25), Inches(top + height / 2 - 0.55),
        Inches(0.5), Inches(0.5),
    )
    icon.fill.solid()
    icon.fill.fore_color.rgb = LIGHT_BLUE
    icon.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(left), Inches(top + height - 0.55), Inches(width), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    p.text = label
    p.font.name = "Calibri"
    p.font.size = Pt(11)
    p.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.CENTER


def add_arrow_down(slide, x: float, y: float, h: float = 0.22) -> None:
    arrow = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.DOWN_ARROW, Inches(x), Inches(y), Inches(0.35), Inches(h),
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = GOLD
    arrow.line.fill.background()


def add_arrow_right(slide, x: float, y: float, w: float = 0.45) -> None:
    arrow = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(0.28),
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = GOLD
    arrow.line.fill.background()


def add_layer_box(slide, x: float, y: float, w: float, h: float, title: str, subtitle: str = "", fill=CARD_BG) -> None:
    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = DARK_BLUE
    box.line.width = Pt(1.2)

    tb = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.08), Inches(w - 0.3), Inches(h - 0.1))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Calibri"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.name = "Calibri"
        p2.font.size = Pt(11)
        p2.font.color.rgb = MUTED
        p2.alignment = PP_ALIGN.CENTER


def add_feature_card(slide, x: float, y: float, icon_char: str, title: str, desc: str) -> None:
    w, h = 3.85, 1.55
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = LIGHT_BLUE
    card.line.width = Pt(1.2)

    badge = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + 0.15), Inches(y + 0.25), Inches(0.55), Inches(0.55),
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = GOLD
    badge.line.fill.background()

    ib = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.32), Inches(0.55), Inches(0.4))
    p = ib.text_frame.paragraphs[0]
    p.text = icon_char
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

    tb = slide.shapes.add_textbox(Inches(x + 0.85), Inches(y + 0.18), Inches(w - 1.0), Inches(h - 0.25))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Calibri"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.name = "Calibri"
    p2.font.size = Pt(11)
    p2.font.color.rgb = MUTED


def slide_shell(slide, prs, page: int, title: str | None = None, subtitle: str | None = None) -> None:
    set_bg(slide, WHITE)
    add_top_accent(slide, prs)
    if title:
        add_title(slide, title)
    if subtitle:
        add_subtitle(slide, subtitle, top=0.82)
    add_footer(slide, prs, page)


# ---------------------------------------------------------------------------
# Slide builders (1–15)
# ---------------------------------------------------------------------------
def slide_01_home(slide, prs, page: int) -> None:
    set_bg(slide, WHITE)
    add_top_accent(slide, prs)

    hero = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0.09), prs.slide_width, Inches(4.6),
    )
    hero.fill.solid()
    hero.fill.fore_color.rgb = DARK_BLUE
    hero.line.fill.background()

    # Decorative circles
    for cx, cy, r in [(11.5, 1.2, 1.8), (1.2, 3.5, 1.2)]:
        c = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(cx), Inches(cy), Inches(r), Inches(r))
        c.fill.solid()
        c.fill.fore_color.rgb = LIGHT_BLUE
        c.fill.transparency = 0.75
        c.line.fill.background()

    add_image_placeholder(slide, 0.55, 0.55, 2.4, 2.4, "App Logo")

    tb = slide.shapes.add_textbox(Inches(3.3), Inches(1.2), Inches(9.2), Inches(1.0))
    p = tb.text_frame.paragraphs[0]
    p.text = "Vithry App"
    p.font.name = "Calibri"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE

    sb = slide.shapes.add_textbox(Inches(3.3), Inches(2.25), Inches(9.2), Inches(0.5))
    p = sb.text_frame.paragraphs[0]
    p.text = "Team Vithey"
    p.font.name = "Calibri"
    p.font.size = Pt(26)
    p.font.color.rgb = GOLD

    tg = slide.shapes.add_textbox(Inches(3.3), Inches(2.9), Inches(9.2), Inches(0.45))
    p = tg.text_frame.paragraphs[0]
    p.text = "Organize · Share · Learn · Career"
    p.font.name = "Calibri"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = LIGHT_BLUE

    pill_y = 5.0
    for i, label in enumerate(["Social", "Jobs", "Finance", "Chat", "AI"]):
        x = 0.55 + i * 2.5
        pill = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(pill_y), Inches(2.2), Inches(0.42),
        )
        pill.fill.solid()
        pill.fill.fore_color.rgb = LIGHT_BLUE
        pill.line.fill.background()
        pb = slide.shapes.add_textbox(Inches(x), Inches(pill_y + 0.07), Inches(2.2), Inches(0.3))
        p = pb.text_frame.paragraphs[0]
        p.text = label
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    add_footer(slide, prs, page)


def slide_02_team(slide, prs, page: int) -> None:
    slide_shell(slide, prs, page, "Team Members")

    badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(10.8), Inches(0.2), Inches(2.0), Inches(0.55))
    badge.fill.solid()
    badge.fill.fore_color.rgb = GOLD
    badge.line.fill.background()
    bb = slide.shapes.add_textbox(Inches(10.8), Inches(0.32), Inches(2.0), Inches(0.3))
    p = bb.text_frame.paragraphs[0]
    p.text = "7 Members"
    p.font.bold = True
    p.font.size = Pt(13)
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

    y = 1.15
    for i, name in enumerate(TEAM):
        row_y = y + i * 0.72
        bg = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(row_y), Inches(10.3), Inches(0.58),
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = CARD_BG if i % 2 == 0 else WHITE
        bg.line.color.rgb = LIGHT_BLUE

        dot = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL, Inches(1.75), Inches(row_y + 0.14), Inches(0.3), Inches(0.3),
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = GOLD
        dot.line.fill.background()

        nb = slide.shapes.add_textbox(Inches(2.2), Inches(row_y + 0.12), Inches(9), Inches(0.35))
        p = nb.text_frame.paragraphs[0]
        p.text = name
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE

    add_image_placeholder(slide, 0.55, 1.15, 0.75, 5.0, "Team Photo")


def slide_03_toc(slide, prs, page: int) -> None:
    slide_shell(slide, prs, page, "Table of Contents")

    cols = 2
    per_col = (len(TOC_ITEMS) + 1) // cols
    for i, item in enumerate(TOC_ITEMS):
        col = i // per_col
        row = i % per_col
        x = 0.7 + col * 6.3
        y = 1.2 + row * 0.62
        num = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y + 0.05), Inches(0.38), Inches(0.38),
        )
        num.fill.solid()
        num.fill.fore_color.rgb = DARK_BLUE
        num.line.fill.background()
        nb = slide.shapes.add_textbox(Inches(x), Inches(y + 0.1), Inches(0.38), Inches(0.28))
        p = nb.text_frame.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        tb = slide.shapes.add_textbox(Inches(x + 0.5), Inches(y + 0.08), Inches(5.5), Inches(0.35))
        p = tb.text_frame.paragraphs[0]
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_TEXT


def slide_04_overview(slide, prs, page: int) -> None:
    slide_shell(slide, prs, page, "Project Overview", "Unified Student Super App")

    features = [
        ("📱", "Social Feed", "Posts, likes, comments & share"),
        ("💼", "Jobs", "Browse jobs & upload CV to apply"),
        ("💰", "Finance", "Student fees, invoices & payments"),
        ("💬", "Chat", "Private messaging & safety tools"),
        ("🤖", "AI Chatbot", "Vithey assistant for student help"),
        ("🔔", "Notifications", "Smart alerts & deep linking"),
    ]
    positions = [(0.55, 1.35), (4.55, 1.35), (8.55, 1.35), (0.55, 3.15), (4.55, 3.15), (8.55, 3.15)]
    for (icon, title, desc), (x, y) in zip(features, positions):
        add_feature_card(slide, x, y, icon, title, desc)

    flow = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(5.05), Inches(12.2), Inches(0.85),
    )
    flow.fill.solid()
    flow.fill.fore_color.rgb = RGBColor(0xFF, 0xF8, 0xE8)
    flow.line.color.rgb = GOLD
    fb = slide.shapes.add_textbox(Inches(0.75), Inches(5.2), Inches(11.8), Inches(0.55))
    p = fb.text_frame.paragraphs[0]
    p.text = "User Flow:  Login  →  Home Feed  →  Jobs  →  Apply CV  →  Chat  →  Finance  →  AI Help"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER


def slide_05_architecture(slide, prs, page: int) -> None:
    slide_shell(slide, prs, page, "System Architecture")

    layers = [
        ("Flutter App", "Mobile client — iOS & Android"),
        ("Spring Boot REST API", "Microservices + API Gateway"),
        ("Spring Security JWT", "Authentication & authorization"),
        ("Database", "PostgreSQL / MySQL"),
        ("Cloud / Kubernetes", "Scalable cloud deployment"),
    ]
    x, w, h = 3.8, 5.7, 0.72
    y = 1.15
    for i, (title, sub) in enumerate(layers):
        fill = CARD_BG if i % 2 == 0 else RGBColor(0xFF, 0xF8, 0xE8)
        add_layer_box(slide, x, y, w, h, title, sub, fill=fill)
        if i < len(layers) - 1:
            add_arrow_down(slide, x + w / 2 - 0.17, y + h + 0.02, 0.2)
        y += h + 0.28

    add_image_placeholder(slide, 0.45, 1.15, 3.0, 5.2, "Architecture Diagram")


def slide_06_frontend(slide, prs, page: int) -> None:
    slide_shell(slide, prs, page, "Frontend Architecture", "Flutter · GetX · Modular Design")

    for x in (0.45, 6.65):
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.2), Inches(5.9), Inches(4.9),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_BG
        card.line.color.rgb = LIGHT_BLUE

    left_items = [
        "GetX — state management & routing",
        "Dio + HTTPS — REST networking",
        "shared_preferences — local settings",
        "flutter_secure_storage — JWT tokens",
        "encrypt / crypto — data security",
    ]
    right_items = [
        "cached_network_image — media cache",
        "intl — date & number formatting",
        "file_picker / image_picker — uploads",
        "Modular: auth, feed, chat, finance…",
        "11 feature modules implemented",
    ]
    add_bullets(slide, left_items, 0.55, 1.35, 5.8, 15)
    add_bullets(slide, right_items, 6.75, 1.35, 5.8, 15)
    add_image_placeholder(slide, 9.0, 4.5, 3.5, 2.0, "App Screenshot")


def slide_07_backend(slide, prs, page: int) -> None:
    slide_shell(slide, prs, page, "Backend Architecture", "Spring Boot 3+ · Java 21 · Maven")

    blocks = [
        ("API Layer", ["REST APIs", "Swagger / OpenAPI", "API Gateway routing"]),
        ("Security", ["Spring Security", "JWT authentication", "Role-based access"]),
        ("Data Layer", ["Spring Data JPA", "Repository pattern", "PostgreSQL / MySQL"]),
        ("Design", ["Service layer", "DTO mapping", "Exception handling"]),
    ]
    for i, (title, items) in enumerate(blocks):
        col = i % 2
        row = i // 2
        x = 0.55 + col * 6.4
        y = 1.25 + row * 2.55
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(5.9), Inches(2.2),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = DARK_BLUE

        hb = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.12), Inches(5.5), Inches(0.35))
        p = hb.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE

        add_bullets(slide, items, x + 0.25, y + 0.5, 5.4, 14)


def slide_08_infrastructure(slide, prs, page: int) -> None:
    slide_shell(slide, prs, page, "Infrastructure")

    components = [
        ("Flutter Client", "Mobile app"),
        ("Spring Boot API", "Microservices"),
        ("Database Server", "PostgreSQL"),
        ("Auth Layer", "JWT tokens"),
        ("Cloud Deploy", "Docker + K8s"),
    ]
    cx, cy = 6.65, 3.6
    radius = 2.4
    import math
    for i, (title, sub) in enumerate(components):
        angle = math.radians(90 + i * (360 / len(components)))
        x = cx + radius * math.cos(angle) - 1.1
        y = cy + radius * math.sin(angle) - 0.45
        add_layer_box(slide, x, y, 2.2, 0.9, title, sub, fill=WHITE)

    hub = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, Inches(cx - 0.65), Inches(cy - 0.65), Inches(1.3), Inches(1.3),
    )
    hub.fill.solid()
    hub.fill.fore_color.rgb = DARK_BLUE
    hub.line.color.rgb = GOLD
    hb = slide.shapes.add_textbox(Inches(cx - 0.65), Inches(cy - 0.35), Inches(1.3), Inches(0.7))
    p = hb.text_frame.paragraphs[0]
    p.text = "Vithey\nPlatform"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    add_image_placeholder(slide, 0.45, 1.15, 3.5, 2.5, "Cloud / Server")


def slide_09_project_plan(slide, prs, page: int) -> None:
    slide_shell(slide, prs, page, "Project Plan", "Development Timeline")

    phases = [
        "Requirements",
        "UI/UX Design",
        "Backend Dev",
        "Frontend Dev",
        "Integration",
        "Testing",
        "Deployment",
    ]
    x_start, step = 0.55, 1.72
    y = 2.8
    for i, phase in enumerate(phases):
        x = x_start + i * step
        circle = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(0.55), Inches(0.55),
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = GOLD if i < 6 else LIGHT_BLUE
        circle.line.color.rgb = DARK_BLUE
        nb = slide.shapes.add_textbox(Inches(x), Inches(y + 0.1), Inches(0.55), Inches(0.35))
        p = nb.text_frame.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE
        p.alignment = PP_ALIGN.CENTER

        if i < len(phases) - 1:
            line = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT, Inches(x + 0.55), Inches(y + 0.27), Inches(x + step), Inches(y + 0.27),
            )
            line.line.color.rgb = LIGHT_BLUE
            line.line.width = Pt(2.5)
            add_arrow_right(slide, x + step - 0.35, y + 0.12, 0.3)

        lb = slide.shapes.add_textbox(Inches(x - 0.35), Inches(y + 0.65), Inches(1.5), Inches(0.55))
        p = lb.text_frame.paragraphs[0]
        p.text = phase
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE
        p.alignment = PP_ALIGN.CENTER

    legend = slide.shapes.add_textbox(Inches(0.55), Inches(4.5), Inches(12), Inches(1.5))
    tf = legend.text_frame
    tf.word_wrap = True
    items = [
        "✅ Phases 1–6 completed (analysis → testing)",
        "🔲 Phase 7 — production deployment to cloud",
        "Tools: Figma · Flutter · Spring Boot · Docker · GitHub Actions",
    ]
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(15)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(6)


def slide_10_devops(slide, prs, page: int) -> None:
    slide_shell(slide, prs, page, "DevOps Strategy", "CI/CD Pipeline")

    steps = ["Git", "Docker", "CI/CD", "Testing", "Cloud Deploy"]
    x_start, step_w = 0.7, 2.35
    y = 2.5
    for i, label in enumerate(steps):
        x = x_start + i * step_w
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.0), Inches(0.85),
        )
        box.fill.solid()
        box.fill.fore_color.rgb = DARK_BLUE if i % 2 == 0 else LIGHT_BLUE
        box.line.fill.background()
        tb = slide.shapes.add_textbox(Inches(x), Inches(y + 0.22), Inches(2.0), Inches(0.4))
        p = tb.text_frame.paragraphs[0]
        p.text = label
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        if i < len(steps) - 1:
            add_arrow_right(slide, x + 2.05, y + 0.28, 0.28)

    add_bullets(slide, [
        "Git version control — branch-based workflow",
        "Docker — per-service container images",
        "GitHub Actions — automated CI per service",
        "Bruno / Postman — API integration testing",
        "Deploy to GCP · Kubernetes · OpenShift",
    ], 0.55, 3.8, 12, 15)


def slide_11_security(slide, prs, page: int) -> None:
    slide_shell(slide, prs, page, "Security Design")

    chain = [
        "Frontend\nSecure Storage",
        "JWT\nTokens",
        "Spring\nSecurity",
        "RBAC",
        "Secure\nREST API",
    ]
    x_start, step = 0.55, 2.45
    y = 2.2
    for i, label in enumerate(chain):
        x = x_start + i * step
        shield = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.PENTAGON, Inches(x), Inches(y), Inches(2.1), Inches(1.1),
        )
        shield.fill.solid()
        shield.fill.fore_color.rgb = CARD_BG if i % 2 == 0 else RGBColor(0xFF, 0xF8, 0xE8)
        shield.line.color.rgb = DARK_BLUE
        tb = slide.shapes.add_textbox(Inches(x), Inches(y + 0.2), Inches(2.1), Inches(0.7))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE
        p.alignment = PP_ALIGN.CENTER
        if i < len(chain) - 1:
            add_arrow_right(slide, x + 2.15, y + 0.4, 0.28)

    add_bullets(slide, [
        "Frontend: flutter_secure_storage for JWT, encrypted prefs",
        "Backend: BCrypt passwords, gateway token validation",
        "System: role-based access, HTTPS-only communication",
    ], 0.55, 3.7, 12, 15)


def slide_12_deployment(slide, prs, page: int) -> None:
    slide_shell(slide, prs, page, "Deployment Strategy")

    platforms = [
        ("Google Cloud", "Primary cloud target"),
        ("Kubernetes", "Container orchestration"),
        ("OpenShift", "Enterprise deployment"),
        ("Oracle WebLogic", "Legacy compatibility"),
    ]
    for i, (name, desc) in enumerate(platforms):
        x = 0.55 + (i % 2) * 6.4
        y = 1.25 + (i // 2) * 2.1
        add_layer_box(slide, x, y, 5.9, 1.6, name, desc, fill=WHITE)

    add_bullets(slide, [
        "Scalability — horizontal scaling via microservices",
        "High availability — stateless service design",
        "Cloud-native — Eureka discovery + Config Server",
    ], 0.55, 5.5, 12, 14)


def slide_13_impact(slide, prs, page: int) -> None:
    slide_shell(slide, prs, page, "Impact & Benefits")

    benefits = [
        ("🎯", "Centralized Platform", "One app for all student needs"),
        ("📢", "Better Communication", "Campus social networking"),
        ("💼", "Career Growth", "Jobs & CV in one place"),
        ("📊", "Finance Clarity", "Track fees & invoices easily"),
        ("🤖", "AI Assistance", "24/7 student support"),
    ]
    for i, (icon, title, desc) in enumerate(benefits):
        x = 0.55 + (i % 3) * 4.15
        y = 1.2 + (i // 3) * 2.0
        add_feature_card(slide, x, y, icon, title, desc)

    add_image_placeholder(slide, 9.0, 4.2, 3.5, 2.2, "Student Impact")


def slide_14_conclusion(slide, prs, page: int) -> None:
    slide_shell(slide, prs, page, "Conclusion")

    points = [
        "Full-stack scalable student super app",
        "Secure JWT architecture end-to-end",
        "10 microservices + Flutter mobile client",
        "Docker & CI/CD deployment ready",
        "Built for real AUB student daily life",
    ]
    for i, point in enumerate(points):
        y = 1.35 + i * 0.85
        check = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL, Inches(1.0), Inches(y), Inches(0.42), Inches(0.42),
        )
        check.fill.solid()
        check.fill.fore_color.rgb = GOLD
        check.line.fill.background()
        cb = slide.shapes.add_textbox(Inches(1.0), Inches(y + 0.05), Inches(0.42), Inches(0.32))
        p = cb.text_frame.paragraphs[0]
        p.text = "✓"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE
        p.alignment = PP_ALIGN.CENTER

        tb = slide.shapes.add_textbox(Inches(1.65), Inches(y + 0.05), Inches(10), Inches(0.4))
        p = tb.text_frame.paragraphs[0]
        p.text = point
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_TEXT

    add_image_placeholder(slide, 8.5, 1.2, 4.0, 4.5, "App Preview")


def slide_15_thank_you(slide, prs, page: int) -> None:
    set_bg(slide, DARK_BLUE)
    add_top_accent(slide, prs)

    tb = slide.shapes.add_textbox(Inches(1), Inches(2.3), Inches(11.3), Inches(1.0))
    p = tb.text_frame.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    sb = slide.shapes.add_textbox(Inches(1), Inches(3.4), Inches(11.3), Inches(0.55))
    p = sb.text_frame.paragraphs[0]
    p.text = "Team Vithey"
    p.font.size = Pt(30)
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER

    tg = slide.shapes.add_textbox(Inches(1), Inches(4.1), Inches(11.3), Inches(0.45))
    p = tg.text_frame.paragraphs[0]
    p.text = "Vithry App — Organize · Share · Learn · Career"
    p.font.size = Pt(18)
    p.font.color.rgb = LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER

    add_image_placeholder(slide, 5.9, 0.55, 1.5, 1.5, "Logo")
    add_footer(slide, prs, page)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
BUILDERS = [
    slide_01_home,
    slide_02_team,
    slide_03_toc,
    slide_04_overview,
    slide_05_architecture,
    slide_06_frontend,
    slide_07_backend,
    slide_08_infrastructure,
    slide_09_project_plan,
    slide_10_devops,
    slide_11_security,
    slide_12_deployment,
    slide_13_impact,
    slide_14_conclusion,
    slide_15_thank_you,
]


def generate(output: Path = OUTPUT) -> Path:
    prs = new_presentation()
    for i, builder in enumerate(BUILDERS, start=1):
        slide = blank_slide(prs)
        builder(slide, prs, i)
    prs.save(str(output))
    print(f"Presentation saved: {output}")
    print(f"   Slides: {len(BUILDERS)}")
    return output


if __name__ == "__main__":
    generate()
